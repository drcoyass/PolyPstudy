import json
import os
from datetime import datetime
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("⚠️ python-pptx がインストールされていません。インストールを試みます...")
    os.system("python3 -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

def create_monthly_report():
    json_path = "data/latest_papers.json"
    output_path = "output/Monthly_Report.pptx"
    
    # フォルダ準備
    if not os.path.exists("output"): os.makedirs("output")

    # データ読み込み
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    papers = data.get('papers', [])
    stats = {}
    for p in papers:
        for tag in p.get('tags', []):
            stats[tag] = stats.get(tag, 0) + 1

    # PPTX作成
    prs = Presentation()

    # --- スライド1: 表紙 ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ポリリン酸研究 月次インテリジェンス"
    subtitle.text = f"国立PubMedデータベース同期報告書\n生成日: {datetime.now().strftime('%Y年%m月%d日')}\n累計論文数: {len(papers):,}件"

    # --- スライド2: 研究トレンド概要 ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "重点研究領域の分布"
    
    # 表の作成（上位5ジャンル）
    sorted_stats = sorted(stats.items(), key=lambda x: x[1], reverse=True)[:5]
    rows, cols = len(sorted_stats) + 1, 2
    left, top, width, height = Inches(1.5), Inches(2), Inches(6), Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table.cell(0, 0).text = "研究ジャンル"
    table.cell(0, 1).text = "論文数"
    
    for i, (genre, count) in enumerate(sorted_stats):
        table.cell(i+1, 0).text = genre
        table.cell(i+1, 1).text = f"{count:,}件"

    # --- スライド3: インプラント研究の最前線 ---
    implant_papers = [p for p in papers if "インプラント" in p.get('tags', [])][:3]
    if implant_papers:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "最新のインプラント関連研究"
        
        for i, p in enumerate(implant_papers):
            top_pos = Inches(1.5 + (i * 1.8))
            textbox = slide.shapes.add_textbox(Inches(0.5), top_pos, Inches(9), Inches(1.5))
            tf = textbox.text_frame
            tf.word_wrap = True
            
            p_title = tf.add_paragraph()
            p_title.text = f"• {p.get('jp_title', p.get('title'))[:80]}..."
            p_title.font.bold = True
            p_title.font.size = Pt(14)
            
            p_year = tf.add_paragraph()
            p_year.text = f"  [{p.get('year')}] PMID: {p.get('id')}"
            p_year.font.size = Pt(10)

    # --- スライド4: 結語・展望 ---
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "今後の展望"
    content = slide.placeholders[1]
    content.text = ("• 再生医療分野におけるポリリン酸の利用が加速\n"
                    "• 炎症抑制・粘膜バリア強化に関する最新知見の増加\n"
                    "• 臨床応用（歯科インプラント・歯周病治療）への高い期待")

    # 保存
    prs.save(output_path)
    print(f"✨ 月次レポートを生成しました: {output_path}")

if __name__ == "__main__":
    create_monthly_report()
